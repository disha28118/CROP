from pptx import Presentation

# Create a new presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Project Insights & Strategic Recommendations"
subtitle.text = "Prepared for Julia [Client/Company Name]\nApril 2025\nPrepared by: [Your Name / Your Team Name]"

# Slide 2: Executive Summary
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Executive Summary"

content = slide_2.shapes.placeholders[1]
content.text = "Key Insights:\n- Customer engagement increased by 20% from Q1 to Q2.\n\nKey Recommendations:\n- Focus on customer retention strategies and invest in high-performing product categories."

# Slide 3: Introduction
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Introduction"

content = slide_3.shapes.placeholders[1]
content.text = "Project Overview:\n- The goal of this analysis was to assess the performance of our key product categories and identify areas for strategic growth.\n\nAnalytics Tasks Overview:\n- Analyzed Q1 and Q2 data to understand customer engagement and product performance."

# Slide 4: Key Insights
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Key Insights from Our Analysis"

content = slide_4.shapes.placeholders[1]
content.text = "Insight 1:\n- Customer engagement increased by 20% from Q1 to Q2.\n- Positive response to recent marketing efforts.\n\nInsight 2:\n- Product category A outperformed others, contributing 40% more revenue in Q2."

# Slide 5: Data Visualizations
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Data Visualizations"

content = slide_5.shapes.placeholders[1]
content.text = "Visualization 1: Line chart showing sales trends.\nVisualization 2: Pie chart displaying market share breakdown.\nVisualization 3: Bar chart comparing product performance."

# Slide 6: Detailed Insights and Analysis
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_6.shapes.title
title.text = "In-depth Analysis"

content = slide_6.shapes.placeholders[1]
content.text = "The 20% increase in customer engagement indicates the effectiveness of recent campaigns.\n- Focus on 25-34 age group for future growth.\n- Continue leveraging high-performing product categories for Q3."

# Slide 7: Recommendations & Next Steps
slide_7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Recommendations & Next Steps"

content = slide_7.shapes.placeholders[1]
content.text = "Recommendations:\n- Focus on targeted marketing for the 25-34 age group.\n- Invest in high-performing categories.\n- Implement retention strategies for long-term loyalty.\n\nNext Steps:\n- Initiate follow-up meeting to discuss resource allocation."

# Slide 8: Conclusion
slide_8 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_8.shapes.title
title.text = "Conclusion"

content = slide_8.shapes.placeholders[1]
content.text = "Summary of Key Insights:\n- Strong growth in customer engagement and product performance.\n- Focus on retention and high-performing products for continued growth.\n\nFinal Thought:\n- Strategic alignment with these insights will ensure sustained success."

# Slide 9: Questions & Discussion
slide_9 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_9.shapes.title
title.text = "Questions & Discussion"

content = slide_9.shapes.placeholders[1]
content.text = "We'd be happy to answer any questions or discuss the next steps further."

# Save the presentation as a .pptx file
output_file = "Project_Insights_Report.pptx"
presentation.save(output_file)

print(f"Presentation saved as {output_file}")
